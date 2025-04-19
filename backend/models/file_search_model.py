from pathlib import Path
from docx import Document
from pptx import Presentation
import fitz  # PyMuPDF
import json

def extract_text_from_file(file_path):
    suffix = file_path.suffix.lower()
    if suffix == '.txt':
        return file_path.read_text(encoding='utf-8', errors='ignore')
    elif suffix == '.docx':
        doc = Document(file_path)
        return '\n'.join(p.text for p in doc.paragraphs)
    elif suffix == '.pptx':
        prs = Presentation(file_path)
        return '\n'.join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    elif suffix == '.pdf':
        doc = fitz.open(file_path)
        return '\n'.join([page.get_text() for page in doc])
    return ''

def index_attachments(root_dir='attachments', save_path=None, print_summary=True):
    index = []
    root = Path(root_dir)
    if not root.exists():
        raise FileNotFoundError(f"Attachment root directory not found: {root_dir}")

    for note_dir in root.glob('note_*'):
        if not note_dir.is_dir():
            continue
        note_id = note_dir.name.replace('note_', '')
        for file in note_dir.glob('*.*'):
            if file.suffix.lower() not in ['.txt', '.docx', '.pptx', '.pdf']:
                continue
            content = extract_text_from_file(file)
            if content.strip():
                index.append({
                    'note_id': note_id,
                    'filename': file.name,
                    'content': content.strip()
                })

    if save_path:
        Path(save_path).write_text(json.dumps(index, ensure_ascii=False, indent=2))
        if print_summary:
            print(f'✅ Indexed {len(index)} attachments saved to {save_path}')
    elif print_summary:
        print(f'✅ Indexed {len(index)} attachments (not saved)')

    return index
